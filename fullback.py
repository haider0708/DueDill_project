from fastapi import FastAPI, File, UploadFile,Form
from fastapi.middleware.cors import CORSMiddleware
import streamlit as st
import os
import PyPDF2
import docx
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from langchain_google_genai import ChatGoogleGenerativeAI
import google.generativeai as genai
from fastapi import FastAPI
from pydantic import BaseModel
from typing import List
import streamlit as st
import os
import PyPDF2
import docx
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime
from langchain_google_genai import ChatGoogleGenerativeAI
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.enum.dml import MSO_FILL # <<< AJOUTER CET IMPORT
import re
from fastapi import FastAPI, HTTPException
from fastapi.responses import JSONResponse, StreamingResponse
import cryptocompare
import pandas as pd
import datetime
import time
import matplotlib
matplotlib.use('Agg') 
import matplotlib.pyplot as plt
import io
import numpy as np
from statsmodels.tsa.stattools import adfuller
import pmdarima as pm
from arch import arch_model
from fastapi.middleware.cors import CORSMiddleware
from fastapi import FastAPI
from pydantic import BaseModel
import tensorflow as tf
import pickle
import pandas as pd
import json
import ollama
from typing import Optional
import csv
from typing import List



TEMPLATE_PATH = r"template.pptx"  
LOGO_PATH = r"logo.png" 
API_KEY = "AIzaSyDlaSyddA2W2sxuXzdQ3Dhg4czVnc2TTJI"  
MODEL_NAME = "gemini-2.0-flash"


app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"], 
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)



def fetch_full_coin_history(symbol: str, currency: str = "USD") -> pd.DataFrame:
    end_timestamp = int(time.time())
    data = cryptocompare.get_historical_price_day(symbol, currency=currency, limit=2000, toTs=end_timestamp)

    if not data:
        return pd.DataFrame()
    
    df = pd.DataFrame(data)[['time', 'close']]
    df.drop_duplicates(subset='time', inplace=True)
    df['time'] = pd.to_datetime(df['time'], unit='s')
    df.sort_values('time', inplace=True)
    df.reset_index(drop=True, inplace=True)

    return df

def plot_coin_data(df: pd.DataFrame, symbol: str):
    fig, ax = plt.subplots(figsize=(10, 6))
    ax.plot(df['time'], df['close'], label=f'{symbol} Close Price', color='blue')
    ax.set_title(f'{symbol} Historical Price')
    ax.set_xlabel('Date')
    ax.set_ylabel('Price (USD)')
    ax.tick_params(axis='x', rotation=45)
    ax.grid(True)
    ax.legend()

    # Save plot to buffer
    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format="png")
    plt.close(fig)
    buf.seek(0)
    return buf

def train_arima(df: pd.DataFrame) -> str:
    df = df.copy()
    start_date = pd.Timestamp.today() - pd.Timedelta(days=len(df))
    df.index = pd.date_range(start=start_date, periods=len(df), freq='D')

    cp = df[['close']]
    cp_log = np.log(cp)

    adf_test = adfuller(cp_log)
    is_stationary = adf_test[1] <= 0.05

    if not is_stationary:
        cp_log_diff = cp_log.diff().dropna()
        model = pm.auto_arima(cp_log_diff, stepwise=False, seasonal=False, suppress_warnings=True)
        model.fit(cp_log_diff)
        forecast_diff = model.predict(n_periods=180)
        last_log = cp_log.iloc[-1, 0]
        forecast_log = np.r_[last_log, forecast_diff].cumsum()[1:]
    else:
        model = pm.auto_arima(cp_log, stepwise=False, seasonal=False, suppress_warnings=True)
        model.fit(cp_log)
        forecast_log = model.predict(n_periods=180)

    forecast = np.exp(forecast_log)
    forecast_series = pd.Series(forecast)

    final_value = forecast_series.iloc[-1]
    initial_value = cp['close'].iloc[-1]
    trend = "📈 Upward" if final_value > initial_value else "📉 Downward"
    return trend

def fit_garch(returns_series, p, q):
    model = arch_model(returns_series, vol='Garch', p=p, q=q)
    model_fit = model.fit(disp="off")
    return model_fit.aic, model_fit.bic, model_fit

def train_garch(df):
    df["close_return"] = np.log(df["close"] / df["close"].shift(1))
    df.dropna(inplace=True)
    df["close_return"] = df["close_return"] * 10

    p_values = range(1, 6)
    q_values = range(1, 6)

    results = []
    for p in p_values:
        for q in q_values:
            aic, bic, model_fit = fit_garch(df["close_return"], p, q)
            results.append((p, q, aic, bic, model_fit))

    best_model = min(results, key=lambda x: x[2])
    return best_model[4]

def interpret_garch_model(model_fit):
    summary_table = model_fit.summary().tables[1].data[1:]
    interpretation = []

    alpha = sum(float(row[1]) for row in summary_table if 'alpha' in row[0])
    beta = sum(float(row[1]) for row in summary_table if 'beta' in row[0])

    interpretation.append("📈 **Volatility Analysis**:")

    if alpha + beta < 1:
        interpretation.append("• Volatility is **stable** over time.")
    else:
        interpretation.append("• Volatility tends to **persist** over time.")

    if beta > alpha:
        interpretation.append("• Past volatility has a stronger influence on future volatility.")
    else:
        interpretation.append("• Recent price movements drive future volatility more than past trends.")

    return interpretation

text_vectorizer = tf.keras.layers.TextVectorization(
    split="whitespace",
    max_tokens=5000,
    output_sequence_length=500
)

model = tf.keras.models.load_model("model_experiments/model_1.keras")



class BytecodeRequest(BaseModel):
    bytecode: str

def split_text_into_chars(text, length=2):
    return " ".join([text[i:i+length] for i in range(0, len(text), length)])


llm = ChatGoogleGenerativeAI(
    temperature=0.1,
    model=MODEL_NAME,
    api_key=API_KEY,
)

def format_text_as_bullets(
    text_frame, 
    text_content, 
    font_name="Georgia", 
    font_size=22, 
    color_rgb=RGBColor(0,0,0), 
    line_spacing_val=1.5, 
    space_after_pt=12, 
    space_before_pt=6, 
    is_continuation=False, 
    first_line_spacing=1.5, 
    continuation_line_spacing=1.0,
    force_bullet_first_line=True
):
    text_frame.clear()
    text_frame.word_wrap = True

    lines = text_content.replace('\r\n', '\n').replace('\r', '\n').strip().split('\n')
    has_added_any_text = False
    for i, line_text in enumerate(lines):
        processed_line = line_text.strip()
        cleaned_prefixes = True
        while cleaned_prefixes:
            cleaned_prefixes = False
            if processed_line.startswith("* "):
                processed_line = processed_line[2:].strip()
                cleaned_prefixes = True
            elif processed_line.startswith("- "):
                processed_line = processed_line[2:].strip()
                cleaned_prefixes = True
            elif processed_line.startswith("• "): 
                processed_line = processed_line[2:].strip()
                cleaned_prefixes = True

        if not processed_line:
            continue

        p = text_frame.add_paragraph()
        if i == 0 and is_continuation and not force_bullet_first_line:
            p.text = processed_line
        else:
            p.text = f"• {processed_line}"

        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(space_before_pt if i == 0 else 3)
        p.space_after = Pt(space_after_pt)
        if i == 0 and is_continuation:
            p.line_spacing = continuation_line_spacing
        else:
            p.line_spacing = first_line_spacing

        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = color_rgb
        has_added_any_text = True

    if not has_added_any_text and text_content.strip():
        p = text_frame.add_paragraph()
        p.text = text_content.strip() 
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = color_rgb
from PyPDF2 import PdfReader
from io import BytesIO
async def lire_document(uploaded_file, filename):
    extension = filename.split('.')[-1].lower()
    texte = ''
    
    file_bytes = await uploaded_file.read()
    file_stream = BytesIO(file_bytes)

    if extension == 'pdf':
        pdf_reader = PdfReader(file_stream)
        for page in pdf_reader.pages:
            if page.extract_text():
                texte += page.extract_text()
    elif extension == 'docx':
        doc = docx.Document(file_stream)
        texte = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    else:
        raise ValueError("Format de fichier non supporté (PDF ou DOCX uniquement).")
    
    return texte


def analyser_document(texte, user_questions, langue="Français"):
    documents = []
    current_doc = ""
    doc_name = ""  
    
    for line in texte.split('\n'):
        if line.startswith('--- DOCUMENT:'):
            if current_doc:
                documents.append((doc_name, current_doc.strip()))
            doc_name = line.split('--- DOCUMENT:')[1].split('---')[0].strip()
            current_doc = ""
        else:
            current_doc += line + '\n'
    if current_doc:
        documents.append((doc_name, current_doc.strip()))

    if not documents:
        st.warning("Aucun document n'a été correctement extrait. Vérifiez le format des fichiers.")
        return [], []

    resumes = []
    for doc_idx, (doc_name, doc_text) in enumerate(documents, 1):
        if not doc_text or len(doc_text) < 50: 
            continue
            
        if langue == "Français":
            resume_prompt = f"""
Document {doc_name} :
\"\"\"{doc_text[:8000]}\"\"\"
Fais un résumé clair en 150 mots maximum en FRANÇAIS. 
"""
        else:
            resume_prompt = f"""
Document {doc_name}:
\"\"\"{doc_text[:8000]}\"\"\"
Provide a clear and structured summary in maximum 150 words in ENGLISH.
"""
        try:
            resume_response = llm.invoke(resume_prompt)
            resume_content = resume_response.content.strip()
            if not resume_content.startswith("Veuillez fournir"):
                resumes.append((doc_name, resume_content))
            else:
                st.warning(f"Impossible de générer un résumé pour le document {doc_name}")
        except Exception as e:
            st.error(f"Erreur lors de la génération du résumé pour {doc_name}: {e}")
            continue

    qa_pairs = []
    for question in user_questions:
        if langue == "Français":
            question_prompt = f"""
Lis attentivement TOUS les documents suivants avant de répondre :
{chr(10).join([f"- {name}: '{doc_text[:200]}...'" for name, doc_text in documents])}

En te basant EXCLUSIVEMENT sur le contenu des documents fournis ci-dessus, réponds de manière précise à la question suivante :
Question : {question}

Instructions pour ta réponse :
Formule ta réponse en FRANÇAIS.
Structure ta réponse sous forme de points avec des puces (•). Chaque point doit correspondre à une idée ou une phrase et être sur une nouvelle ligne.

3.  Indique clairement dans une section distincte nommée "*Référence :*" quel document (avec son titre complet tel que fourni ci-dessus) et quelle(s) phrase(s) spécifique(s) tu as utilisé(s) pour formuler ta réponse.
4.  Pour la section "*Référence :*", fournis une liste simple. Chaque élément de la liste doit commencer UNIQUEMENT par une puce (•) suivie d'un espace, puis du texte de la référence. N'utilise AUCUN autre marqueur comme des astérisques (*) ou des tirets (-). Chaque référence doit être sur sa propre ligne.
"""
        else: 
            question_prompt = f"""
Carefully read ALL the following documents before answering:
{chr(10).join([f"- {name}: '{doc_text[:200]}...'" for name, doc_text in documents])}

Based EXCLUSIVELY on the content of the documents provided above, provide a precise answer to the following question:
Question: {question}

Instructions for your answer:
Formulate your answer in ENGLISH.
Structure your answer in bullet points (•). Each point should correspond to an idea or sentence and be on a new line.

3.  Clearly indicate in a separate section named "*Reference:*" which document (with its full title as provided above) and which specific sentence(s) you used to formulate your answer.
4.  For the "*Reference:*" section, provide a simple list. Each list item must start ONLY with a bullet point (•) followed by a space, then the reference text. Do NOT use ANY other markers like asterisks (*) or hyphens (-). Each reference should be on its own line.
"""
        answer_response = llm.invoke(question_prompt)
        answer = answer_response.content.strip()
        qa_pairs.append((question, answer))

    return resumes, qa_pairs


def get_font_size(text, max_font=36, min_font=16):
    length = len(text)
    if length < 200:
        return max_font
    elif length < 500:
        return int(max_font * 0.75)
    elif length < 1000:
        return int(max_font * 0.5)
    else:
        return min_font

def add_background_image(slide, prs, image_path):
    if os.path.exists(image_path):
        left = top = Inches(0)
        slide.shapes.add_picture(
            image_path,
            left,
            top,
            width=prs.slide_width,
            height=prs.slide_height
        )
def add_slide_border(slide, color, thickness=Pt(1.5)):
    for shape in slide.shapes:
        if shape.line:
            shape.line.color.rgb = color
            shape.line.width = thickness
            shape.line.dash_style = MSO_LINE.SOLID
def add_logo(slide, logo_path, left, top, width):
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, left, top, width=width)

def apply_text_style(text_frame, font_name="Arial", font_size=24, bold=False, color=RGBColor(0, 0, 0), alignment=PP_ALIGN.LEFT):
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color

def add_page_number(slide, page_num_text):
    left = Inches(9.2) 
    top = Inches(7.0) 
    width = Inches(0.5)
    height = Inches(0.4) 
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.clear() 
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = page_num_text
    
    font = run.font
    font.name = 'Arial'
    font.size = Pt(20)
    font.bold = False
    font.color.rgb = RGBColor(0, 0, 0) 
    p.alignment = PP_ALIGN.RIGHT 
    tf.margin_bottom = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.word_wrap = False

    txBox.line.fill.background() 


def split_long_text(text, max_chars=450):
    parts = []
    current_part = ""
    
    paragraphs = text.split("\n\n")
    
    for paragraph in paragraphs:
        paragraph = paragraph.strip()
        if not paragraph:
            continue
            
        if len(paragraph) > max_chars:
            words = paragraph.split()
            for word in words:
                if len(current_part) + len(word) + 1 <= max_chars:
                    current_part += word + " "
                else:
                    if current_part:
                        parts.append(current_part.strip())
                    current_part = word + " "
        else:
            if len(current_part) + len(paragraph) + 2 <= max_chars:
                current_part += paragraph + "\n\n"
            else:
                if current_part:
                    parts.append(current_part.strip())
                current_part = paragraph + "\n\n"
    
    if current_part:
        parts.append(current_part.strip())
    
    return parts

def split_text_by_sentences(text, max_chars=300):
    """
    Découpe un texte en morceaux de phrases, chaque morceau ne dépassant pas max_chars caractères.
    On essaie de couper aux fins de phrases (., !, ?).
    """
    import re
    sentences = re.split(r'(?<=[.!?])\s+', text)
    parts = []
    current_part = ""
    for sentence in sentences:
        if len(current_part) + len(sentence) + 1 <= max_chars:
            current_part += sentence + " "
        else:
            if current_part:
                parts.append(current_part.strip())
            current_part = sentence + " "
    if current_part:
        parts.append(current_part.strip())
    return parts


def load_top_questions(csv_path: str, top_n: int = 5) -> List[str]:
    questions = []
    with open(csv_path, mode='r', encoding='utf-8') as f:
        reader = csv.reader(f)
        next(reader)  # Skip header
        for row in reader:
            if row:
                questions.append(row[0])
            if len(questions) >= top_n:
                break
    return questions

def creer_presentation(resumes, qa_pairs, titre_doc):
    with open(TEMPLATE_PATH, "rb") as f:
        prs = Presentation(f)
    import datetime
    today = datetime.datetime.now().strftime("%d/%m/%Y")
    
    for i in range(len(prs.slides) - 1, -1, -1):
        r_id = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[i])
    
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    
    slide_titre = prs.slides.add_slide(title_slide_layout)
    add_background_image(slide_titre, prs, 'background_title.jpg')
    add_slide_border(slide_titre, RGBColor(255, 215, 0), Pt(3.5))
    
    logo_added = add_logo(slide_titre, LOGO_PATH, Inches(8), Inches(0.5), Inches(1.5))
    title_shape = slide_titre.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(2))
    title_frame = title_shape.text_frame
    title_frame.word_wrap = True
    title_frame.text = f"Document Analysis\n{titre_doc}\n{today}"
    apply_text_style(title_frame, font_name="Georgia", font_size=25,
                     bold=True, color=RGBColor(0, 0, 0), alignment=PP_ALIGN.CENTER)
   
    
    for doc_idx, (doc_name, resume) in enumerate(resumes, 1):
        resume_parts = split_long_text(resume, max_chars=350) 
        for part_idx, resume_part in enumerate(resume_parts, 1):
            slide_resume = prs.slides.add_slide(content_slide_layout)
            add_background_image(slide_resume, prs, 'background_type.jpg')
            add_page_number(slide_resume, str(len(prs.slides)))
            add_slide_border(slide_resume, RGBColor(0, 32, 96))
            
            title_shape = slide_resume.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
            title_frame = title_shape.text_frame
            title_frame.word_wrap = True
            title_frame.text = f"Document Summary {doc_idx}"
            if len(resume_parts) > 1:
                title_frame.text += f" (Part {part_idx}/{len(resume_parts)})"
            apply_text_style(title_frame, font_name="Georgia", font_size=28,
                          bold=True, color=RGBColor(0, 32, 96), alignment=PP_ALIGN.CENTER)
            
            subtitle_shape = slide_resume.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.5))
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.text = doc_name
            apply_text_style(subtitle_frame, font_name="Georgia", font_size=20,
                          bold=False, color=RGBColor(0, 32, 96), alignment=PP_ALIGN.CENTER)
            
            content_shape = slide_resume.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(3.5)) 
            content_frame = content_shape.text_frame
            content_frame.clear() 
            content_frame.word_wrap = True
            content_frame.text = resume_part
            for paragraph in content_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.JUSTIFY 
                paragraph.line_spacing = 1.5
                paragraph.space_before = Pt(5)
                paragraph.space_after = Pt(10)
                for run in paragraph.runs:
                    run.font.name = "Georgia"
                    run.font.size = Pt(20)
                    run.font.color.rgb = RGBColor(0, 0, 0)
    
    for idx, (question, answer) in enumerate(qa_pairs, start=1):
        q_slide = prs.slides.add_slide(content_slide_layout)
        add_background_image(q_slide, prs, 'background_type.jpg')
        add_page_number(q_slide, str(len(prs.slides))) 
        add_slide_border(q_slide, RGBColor(0, 0, 128))
        
        title_shape = q_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = f"Question {idx}"
        apply_text_style(title_frame, font_name="Georgia", font_size=32, 
                         bold=True, color=RGBColor(0, 32, 96), alignment=PP_ALIGN.CENTER)
        
        content_shape = q_slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(3.5)) 
        content_frame = content_shape.text_frame
        content_frame.text = question
        content_frame.word_wrap = True
        for paragraph in content_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER 
            for run in paragraph.runs:
                run.font.name = "Georgia"
                run.font.size = Pt(24) 
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 0, 128)
        
        if "**Référence :**" in answer:
            answer_text, ref_text = answer.split("**Référence :**", 1)
            answer_text = answer_text.strip()
            ref_text = ref_text.strip()
        else:
            answer_text = answer
            ref_text = ""
        
        # Pour les réponses
        answer_parts = split_text_by_sentences(answer_text, max_chars=300)
        last_part_idx = len(answer_parts)
        if not answer_parts and answer_text:
            answer_parts = [answer_text]
            last_part_idx = 1

        for part_idx, answer_part in enumerate(answer_parts, 1):
            r_slide = prs.slides.add_slide(content_slide_layout)
            add_background_image(r_slide, prs, 'background_type.jpg')
            add_page_number(r_slide, str(len(prs.slides)))
            add_slide_border(r_slide, RGBColor(0, 64, 128))
            
            title_shape = r_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = f"Response to question {idx}"
            if len(answer_parts) > 1:
                title_frame.text += f" (Part {part_idx}/{len(answer_parts)})"
            apply_text_style(title_frame, font_name="Georgia", font_size=32, 
                          bold=True, color=RGBColor(0, 64, 128), alignment=PP_ALIGN.CENTER)
            
            is_continuation = part_idx > 1 and not answer_part.strip().startswith("•")
            format_text_as_bullets(
                r_slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(3.5)).text_frame,
                answer_part,
                font_size=20,
                line_spacing_val=1.5,
                space_after_pt=10,
                space_before_pt=5,
                is_continuation=is_continuation,
                force_bullet_first_line=not is_continuation
            )

            if ref_text and part_idx == last_part_idx:
                ref_parts = split_text_by_sentences(ref_text, max_chars=300)
                if not ref_parts and ref_text:
                    ref_parts = [ref_text]

                for ref_part_idx_local, ref_part_content in enumerate(ref_parts, 1):
                    ref_slide = prs.slides.add_slide(content_slide_layout)
                    add_background_image(ref_slide, prs, 'background_type.jpg')
                    add_page_number(ref_slide, str(len(prs.slides)))
                    add_slide_border(ref_slide, RGBColor(0, 64, 128))
                    
                    title_shape_ref = ref_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                    title_frame_ref = title_shape_ref.text_frame
                    title_frame_ref.text = f"References {idx}"
                    if len(ref_parts) > 1:
                        title_frame_ref.text += f" (Partie {ref_part_idx_local}/{len(ref_parts)})"
                    apply_text_style(title_frame_ref, font_name="Georgia", font_size=28,
                                  bold=True, color=RGBColor(0, 64, 128), alignment=PP_ALIGN.CENTER)
                    
                    format_text_as_bullets(
                        ref_slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(3.5)).text_frame,
                        ref_part_content,
                        font_size=18,
                        line_spacing_val=1.2,
                        space_after_pt=5,
                        space_before_pt=3,
                        is_continuation=False,
                        force_bullet_first_line=True
                    )

    safe_filename_component = titre_doc.replace(' ', '_')
    safe_filename_component = re.sub(r'[^\w\._-]', '_', safe_filename_component)
    safe_filename_component = re.sub(r'__+', '_', safe_filename_component)
    safe_filename_component = safe_filename_component.strip('_')
    if not safe_filename_component:
        safe_filename_component = "document"

    output_path = f"Analyse_Document_{safe_filename_component}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(output_path)
    return output_path


with open("crypto_listing.json", "r") as f:
    listings_data = json.load(f)["data"]  
news_data = pd.read_csv("news_data.csv")

def get_crypto_data(symbol, listings_data, news_df):
    crypto_info = next((item for item in listings_data if item['symbol'].lower() == symbol.lower()), None)
    if not crypto_info:
        return None, None
    news_related = news_df[news_df['trading_pair_or_currencies'].str.lower().str.contains(symbol.lower(), na=False)].sort_values(by="time", ascending=False).head(5)

    return crypto_info, news_related.to_dict(orient='records')


def ask_llama_about_risk(symbol, market_info, news_related):
    prompt = f"""
You are a financial risk analyst. A user is asking about the RISK of investing in the cryptocurrency '{symbol}'.

Here is the market data:
{json.dumps(market_info, indent=2)}

Here are recent news articles:
{json.dumps(news_related, indent=2)}

Based on this information, tell me:
- Is this cryptocurrency risky?
- Why or why not?
- Summarize the decision in a single line at the end: RISK LEVEL = [LOW/MEDIUM/HIGH]

Limit your response to **no more than 4 to 5 concise lines**.
Don't forget to give The latest value of the coin.
"""

    response = ollama.chat(
        model='llama3.2',
        messages=[{"role": "user", "content": prompt}]
    )
    return response['message']['content']


@app.post("/predict")
def predict(request: BytecodeRequest):
    bytecode_processed = split_text_into_chars(request.bytecode)
    bytecode_tensor = tf.constant([[bytecode_processed]])
    predictions_probs = model.predict(bytecode_tensor)
    preds = {}
    for i, prob in enumerate(predictions_probs):
        preds[str(i)] = int(prob[0] > 0.5)  # Threshold 0.5
    return {"predictions": preds}


@app.get("/plot/{symbol}")
def get_plot(symbol: str):
    df = fetch_full_coin_history(symbol)
    if df.empty:
        raise HTTPException(status_code=404, detail="No data found for this symbol.")
    buf = plot_coin_data(df, symbol)
    return StreamingResponse(buf, media_type="image/png")

@app.get("/forecast/{symbol}")
def get_forecast(symbol: str):
    df = fetch_full_coin_history(symbol)
    if df.empty:
        raise HTTPException(status_code=404, detail="No data found for this symbol.")
    
    try:
        trend = train_arima(df)
        model = train_garch(df)
        interpretation = interpret_garch_model(model)
        return JSONResponse(content={
            "symbol": symbol,
            "trend": trend,
            "volatility_analysis": interpretation
        })
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/upload")
async def upload_file(
    questions: List[str],
    files: List[UploadFile] = File(...)
):
    questions = questions[0].split(',')
    all_text = ""
    file_titles = []
    for file in files:
                    file_text = await lire_document(file,file.filename)
                    all_text += f"\n\n--- DOCUMENT: {file.filename} ---\n\n{file_text}"
                    file_titles.append(file.filename)
    print(questions)
    resume, qa_pairs = analyser_document(
                        all_text, 
                        questions,
                        "English"
                    )
    for i, (question, reponse) in enumerate(qa_pairs, 1):
        extrait = reponse[:150] + ("..." if len(reponse) > 150 else "")
    titre_doc = " + ".join([os.path.splitext(name)[0] for name in file_titles])               
    pptx_path = creer_presentation(resume, qa_pairs, titre_doc)
    if pptx_path and os.path.exists(pptx_path):
        st.success("✅ Présentation générée avec succès !")
        with open(pptx_path, "rb") as f:
            st.download_button(
            label="⬇️ Télécharger le PowerPoint",
            data=f,
            file_name=os.path.basename(pptx_path), 
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
    return (pptx_path)



@app.get("/news/{query}")
def get_news(query: Optional[str] = None):
    market_info, news_related = get_crypto_data(query, listings_data, news_data)
    if market_info is None:
            risk_level = "No data available for the provided symbol."
    else:
            risk_level = ask_llama_about_risk(query, market_info, news_related)

    return {risk_level}

@app.post("/upload5")
async def upload_file(files: List[UploadFile] = File(...)):
    csv_path = "extracted_questions.csv"  # Your saved question bank CSV
    questions = load_top_questions(csv_path, top_n=5)

    all_text = ""
    file_titles = []
    for file in files:
        file_text = await lire_document(file, file.filename)
        all_text += f"\n\n--- DOCUMENT: {file.filename} ---\n\n{file_text}"
        file_titles.append(file.filename)

    print("Using Questions:", questions)

    resume, qa_pairs = analyser_document(
        all_text, 
        questions,
        "English"
    )

    for i, (question, reponse) in enumerate(qa_pairs, 1):
        extrait = reponse[:150] + ("..." if len(reponse) > 150 else "")

    titre_doc = " + ".join([os.path.splitext(name)[0] for name in file_titles])               
    pptx_path = creer_presentation(resume, qa_pairs, titre_doc)

    return {"pptx_path": pptx_path}

