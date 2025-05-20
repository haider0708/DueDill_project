# --- IMPORTS ---
import streamlit as st
import os
import PyPDF2
import docx
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime
from langchain_google_genai import ChatGoogleGenerativeAI
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.enum.dml import MSO_FILL # <<< AJOUTER CET IMPORT
import re

# --- CONFIG ---
TEMPLATE_PATH = r"C:\Users\Mega Pc\Desktop\Nouveau dossier\template.pptx"
LOGO_PATH = r"logo.jpg"  
API_KEY = "AIzaSyDlaSyddA2W2sxuXzdQ3Dhg4czVnc2TTJI"  
MODEL_NAME = "gemini-2.0-flash"

# --- INITIALISATION LLM ---
llm = ChatGoogleGenerativeAI(
    temperature=0.1,
    model=MODEL_NAME,
    api_key=API_KEY,
)

def format_text_as_bullets(
    text_frame, 
    text_content, 
    font_name="Georgia", 
    font_size=22, 
    color_rgb=RGBColor(0,0,0), 
    line_spacing_val=1.5, 
    space_after_pt=12, 
    space_before_pt=6, 
    is_continuation=False, 
    first_line_spacing=1.5, 
    continuation_line_spacing=1.0,
    force_bullet_first_line=True
):
    text_frame.clear()
    text_frame.word_wrap = True

    lines = text_content.replace('\r\n', '\n').replace('\r', '\n').strip().split('\n')
    has_added_any_text = False
    for i, line_text in enumerate(lines):
        processed_line = line_text.strip()
        # Nettoyage des préfixes
        cleaned_prefixes = True
        while cleaned_prefixes:
            cleaned_prefixes = False
            if processed_line.startswith("* "):
                processed_line = processed_line[2:].strip()
                cleaned_prefixes = True
            elif processed_line.startswith("- "):
                processed_line = processed_line[2:].strip()
                cleaned_prefixes = True
            elif processed_line.startswith("• "): 
                processed_line = processed_line[2:].strip()
                cleaned_prefixes = True

        if not processed_line:
            continue

        p = text_frame.add_paragraph()
        # Si c'est la première ligne et qu'on ne veut pas de puce (suite d'une phrase)
        if i == 0 and is_continuation and not force_bullet_first_line:
            p.text = processed_line
        else:
            p.text = f"• {processed_line}"

        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(space_before_pt if i == 0 else 3)
        p.space_after = Pt(space_after_pt)
        # Interligne dynamique
        if i == 0 and is_continuation:
            p.line_spacing = continuation_line_spacing
        else:
            p.line_spacing = first_line_spacing

        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = color_rgb
        has_added_any_text = True

    if not has_added_any_text and text_content.strip():
        p = text_frame.add_paragraph()
        p.text = text_content.strip() 
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = color_rgb

# --- LIRE DOCUMENT ---
def lire_document(uploaded_file):
    extension = uploaded_file.name.split('.')[-1].lower()
    texte = ''
    if extension == 'pdf':
        pdf_reader = PyPDF2.PdfReader(uploaded_file)
        for page in pdf_reader.pages:
            if page.extract_text():
                texte += page.extract_text()
    elif extension == 'docx':
        doc = docx.Document(uploaded_file)
        texte = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    else:
        raise ValueError("Format de fichier non supporté (PDF ou DOCX uniquement).")
    return texte

# --- ANALYSER DOCUMENT ---
def analyser_document(texte, user_questions, langue="Français"):
    # Extraire les documents individuels
    documents = []
    current_doc = ""
    doc_name = ""  # Initialisation de la variable doc_name
    
    for line in texte.split('\n'):
        if line.startswith('--- DOCUMENT:'):
            if current_doc:
                documents.append((doc_name, current_doc.strip()))
            doc_name = line.split('--- DOCUMENT:')[1].split('---')[0].strip()
            current_doc = ""
        else:
            current_doc += line + '\n'
    if current_doc:
        documents.append((doc_name, current_doc.strip()))

    # Vérifier que les documents ont été correctement extraits
    if not documents:
        st.warning("Aucun document n'a été correctement extrait. Vérifiez le format des fichiers.")
        return [], []

    # Générer les résumés par document
    resumes = []
    for doc_idx, (doc_name, doc_text) in enumerate(documents, 1):
        if not doc_text or len(doc_text) < 50:  # Vérifier que le document contient du texte
            continue
            
        if langue == "Français":
            resume_prompt = f"""
Document {doc_name} :
\"\"\"{doc_text[:8000]}\"\"\"
Fais un résumé clair en 150 mots maximum en FRANÇAIS. 
"""
        else:
            resume_prompt = f"""
Document {doc_name}:
\"\"\"{doc_text[:8000]}\"\"\"
Provide a clear and structured summary in maximum 150 words in ENGLISH.
"""
        try:
            resume_response = llm.invoke(resume_prompt)
            resume_content = resume_response.content.strip()
            # Vérifier que le résumé n'est pas un message d'erreur
            if not resume_content.startswith("Veuillez fournir"):
                resumes.append((doc_name, resume_content))
            else:
                st.warning(f"Impossible de générer un résumé pour le document {doc_name}")
        except Exception as e:
            st.error(f"Erreur lors de la génération du résumé pour {doc_name}: {e}")
            continue

    # Traiter les questions
    qa_pairs = []
    for question in user_questions:
        if langue == "Français":
            question_prompt = f"""
Documents analysés :
{' '.join([f'- {name}' for name, _ in documents])}

Question : {question}
Fournis une réponse précise en FRANÇAIS basée sur ces documents.
Indique clairement dans la section "**Référence :**" quel document (avec son titre) et quelle phrase spécifique tu utilises pour ta réponse.
Pour la section Référence, fournis une liste simple. Chaque élément de la liste doit commencer UNIQUEMENT par une puce (•) suivie d'un espace, puis du texte de la référence. N'utilise AUCUN autre marqueur comme des astérisques (*) ou des tirets (-). Chaque référence doit être sur sa propre ligne.
Les réponses doivent aussi être sous forme de points avec des puces (•) et des retours à la ligne. Pour les références et les réponses il faut un retour a la ligne pour chaque phrase.
"""
        else: # Pour l'anglais, adapter de manière similaire si besoin
            question_prompt = f"""
Analyzed documents:
{' '.join([f'- {name}' for name, _ in documents])}

Question: {question}
Provide a precise answer in ENGLISH based on these documents.
Clearly indicate in the "**Reference:**" section which document (with its title) and specific sentence you're using for your answer.
For the Reference section, provide a simple list. Each list item must start ONLY with a bullet point (•) followed by a space, then the reference text. Do NOT use ANY other markers like asterisks (*) or hyphens (-). Each reference should be on its own line.
Answers should also be in bullet points (•) with line breaks.
"""
        answer_response = llm.invoke(question_prompt)
        answer = answer_response.content.strip()
        qa_pairs.append((question, answer))

    return resumes, qa_pairs

# --- AJUSTEMENT DYNAMIQUE POLICE ---
def get_font_size(text, max_font=36, min_font=16):
    length = len(text)
    if length < 200:
        return max_font
    elif length < 500:
        return int(max_font * 0.75)
    elif length < 1000:
        return int(max_font * 0.5)
    else:
        return min_font

def add_background_image(slide, prs, image_path):
    if os.path.exists(image_path):
        left = top = Inches(0)
        slide.shapes.add_picture(
            image_path,
            left,
            top,
            width=prs.slide_width,
            height=prs.slide_height
        )
def add_slide_border(slide, color, thickness=Pt(1.5)):
    for shape in slide.shapes:
        # Vérifie si la forme a déjà une ligne
        if shape.line:
            shape.line.color.rgb = color
            shape.line.width = thickness
            shape.line.dash_style = MSO_LINE.SOLID
def add_logo(slide, logo_path, left, top, width):
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, left, top, width=width)

def apply_text_style(text_frame, font_name="Arial", font_size=24, bold=False, color=RGBColor(0, 0, 0), alignment=PP_ALIGN.LEFT):
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color

def add_page_number(slide, page_num_text):
    # Position en bas à droite
    # La largeur standard d'une diapositive est de 10 pouces.
    # Nous positionnons la boîte de texte de 0.5 pouces de large à 9.2 pouces du bord gauche.
    left = Inches(9.2) # Ajusté pour le côté droit
    top = Inches(7.0) # Conserver la même position verticale
    width = Inches(0.5) # Conserver la même largeur
    height = Inches(0.4) # Conserver la même hauteur
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.clear() 
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = page_num_text
    
    font = run.font
    font.name = 'Arial'
    font.size = Pt(20)
    font.bold = False
    font.color.rgb = RGBColor(0, 0, 0) 
    p.alignment = PP_ALIGN.RIGHT # Changé pour aligner à droite
    tf.margin_bottom = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.word_wrap = False

    # Supprimer le cadre de la zone de texte
    txBox.line.fill.background() # Définit le remplissage de la ligne sur aucun remplissage


# --- DIVISER TEXTE LONG ---
def split_long_text(text, max_chars=450):
    parts = []
    current_part = ""
    
    # Diviser d'abord par paragraphes naturels
    paragraphs = text.split("\n\n")
    
    for paragraph in paragraphs:
        paragraph = paragraph.strip()
        if not paragraph:
            continue
            
        # Si le paragraphe est trop long, le diviser en morceaux
        if len(paragraph) > max_chars:
            words = paragraph.split()
            for word in words:
                if len(current_part) + len(word) + 1 <= max_chars:
                    current_part += word + " "
                else:
                    if current_part:
                        parts.append(current_part.strip())
                    current_part = word + " "
        else:
            if len(current_part) + len(paragraph) + 2 <= max_chars:
                current_part += paragraph + "\n\n"
            else:
                if current_part:
                    parts.append(current_part.strip())
                current_part = paragraph + "\n\n"
    
    if current_part:
        parts.append(current_part.strip())
    
    return parts

def split_text_by_sentences(text, max_chars=300):
    """
    Découpe un texte en morceaux de phrases, chaque morceau ne dépassant pas max_chars caractères.
    On essaie de couper aux fins de phrases (., !, ?).
    """
    import re
    sentences = re.split(r'(?<=[.!?])\s+', text)
    parts = []
    current_part = ""
    for sentence in sentences:
        if len(current_part) + len(sentence) + 1 <= max_chars:
            current_part += sentence + " "
        else:
            if current_part:
                parts.append(current_part.strip())
            current_part = sentence + " "
    if current_part:
        parts.append(current_part.strip())
    return parts

def creer_presentation(resumes, qa_pairs, titre_doc):
    # Créer une nouvelle présentation à partir du template
    with open(TEMPLATE_PATH, "rb") as f:
        prs = Presentation(f)
    today = datetime.now().strftime("%d/%m/%Y")
    
    # Supprimer les slides existants du template
    for i in range(len(prs.slides) - 1, -1, -1):
        r_id = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[i])
    
    # Obtenir les layouts
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    
    # Add title slide
    slide_titre = prs.slides.add_slide(title_slide_layout)
    add_background_image(slide_titre, prs, 'C:/Users/Lyna Bouzouita/projet/background_title.jpg')
    # Si vous souhaitez un numéro de page sur la diapositive de titre, décommentez la ligne suivante :
    # add_page_number(slide_titre, str(len(prs.slides))) 
    add_slide_border(slide_titre, RGBColor(255, 215, 0), Pt(3.5))
    
    # Ajouter le logo sur la première diapo
    logo_added = add_logo(slide_titre, LOGO_PATH, Inches(8), Inches(0.5), Inches(1.5))
    
    # Add title to the title slide
    title_shape = slide_titre.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(2))
    title_frame = title_shape.text_frame
    title_frame.word_wrap = True
    title_frame.text = f"Analyse de Document\n{titre_doc}\n{today}"
    apply_text_style(title_frame, font_name="Georgia", font_size=25,
                     bold=True, color=RGBColor(0, 0, 0), alignment=PP_ALIGN.CENTER)
    # Pas de numéro de page sur la diapositive de titre, ou commentez la ligne suivante si vous en voulez un.
    # add_page_number(slide_titre, str(len(prs.slides))) 
    
    # 1. Ajouter tous les résumés des documents en premier
    for doc_idx, (doc_name, resume) in enumerate(resumes, 1):
        resume_parts = split_long_text(resume, max_chars=350) 
        for part_idx, resume_part in enumerate(resume_parts, 1):
            slide_resume = prs.slides.add_slide(content_slide_layout)
            # Correction de l'ordre ici :
            add_background_image(slide_resume, prs, 'C:/Users/Lyna Bouzouita/projet/background_type.jpg')
            add_page_number(slide_resume, str(len(prs.slides))) # DÉPLACÉ APRÈS add_background_image
            add_slide_border(slide_resume, RGBColor(0, 32, 96))
            
            title_shape = slide_resume.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
            title_frame = title_shape.text_frame
            title_frame.word_wrap = True
            title_frame.text = f"Résumé du Document {doc_idx}"
            if len(resume_parts) > 1:
                title_frame.text += f" (Partie {part_idx}/{len(resume_parts)})"
            apply_text_style(title_frame, font_name="Georgia", font_size=28,
                          bold=True, color=RGBColor(0, 32, 96), alignment=PP_ALIGN.CENTER)
            
            subtitle_shape = slide_resume.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.5))
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.text = doc_name
            apply_text_style(subtitle_frame, font_name="Georgia", font_size=20,
                          bold=False, color=RGBColor(0, 32, 96), alignment=PP_ALIGN.CENTER)
            
            content_shape = slide_resume.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(3.5)) 
            content_frame = content_shape.text_frame
            content_frame.clear() 
            content_frame.word_wrap = True
            # Les résumés restent en format paragraphe
            content_frame.text = resume_part
            for paragraph in content_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.JUSTIFY 
                paragraph.line_spacing = 1.5
                paragraph.space_before = Pt(5)
                paragraph.space_after = Pt(10)
                for run in paragraph.runs:
                    run.font.name = "Georgia"
                    run.font.size = Pt(20)
                    run.font.color.rgb = RGBColor(0, 0, 0)
    
    # 2. ENSUITE traiter toutes les questions/réponses (APRÈTOUS LES RÉSUMÉS)
    for idx, (question, answer) in enumerate(qa_pairs, start=1):
        q_slide = prs.slides.add_slide(content_slide_layout)
        add_background_image(q_slide, prs, 'C:/Users/Lyna Bouzouita/projet/background_type.jpg')
        add_page_number(q_slide, str(len(prs.slides))) # AJOUT DU NUMÉRO DE PAGE
        add_slide_border(q_slide, RGBColor(0, 0, 128))
        
        title_shape = q_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = f"Question {idx}"
        apply_text_style(title_frame, font_name="Georgia", font_size=32, 
                         bold=True, color=RGBColor(0, 32, 96), alignment=PP_ALIGN.CENTER)
        
        # Ajuster position et hauteur du contenu de la question
        content_shape = q_slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(3.5)) 
        content_frame = content_shape.text_frame
        # La question elle-même n'est pas formatée en puces, mais on peut centrer et styler
        content_frame.text = question
        content_frame.word_wrap = True
        for paragraph in content_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER 
            for run in paragraph.runs:
                run.font.name = "Georgia"
                run.font.size = Pt(24) # Taille ajustée pour la question
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 0, 128)
        
        if "**Référence :**" in answer:
            answer_text, ref_text = answer.split("**Référence :**", 1)
            answer_text = answer_text.strip()
            ref_text = ref_text.strip()
        else:
            answer_text = answer
            ref_text = ""
        
        # Pour les réponses
        answer_parts = split_text_by_sentences(answer_text, max_chars=300)
        last_part_idx = len(answer_parts)
        if not answer_parts and answer_text:
            answer_parts = [answer_text]
            last_part_idx = 1

        for part_idx, answer_part in enumerate(answer_parts, 1):
            r_slide = prs.slides.add_slide(content_slide_layout)
            add_background_image(r_slide, prs, 'C:/Users/Lyna Bouzouita/projet/background_type.jpg')
            add_page_number(r_slide, str(len(prs.slides)))
            add_slide_border(r_slide, RGBColor(0, 64, 128))
            
            title_shape = r_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = f"Réponse à la question {idx}"
            if len(answer_parts) > 1:
                title_frame.text += f" (Partie {part_idx}/{len(answer_parts)})"
            apply_text_style(title_frame, font_name="Georgia", font_size=32, 
                          bold=True, color=RGBColor(0, 64, 128), alignment=PP_ALIGN.CENTER)
            
            # Détection de la suite de phrase (pas de puce si la phrase continue)
            is_continuation = part_idx > 1 and not answer_part.strip().startswith("•")
            format_text_as_bullets(
                r_slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(3.5)).text_frame,
                answer_part,
                font_size=20,
                line_spacing_val=1.5,
                space_after_pt=10,
                space_before_pt=5,
                is_continuation=is_continuation,
                force_bullet_first_line=not is_continuation
            )

            # Références sur une nouvelle slide après la dernière partie de la réponse
            if ref_text and part_idx == last_part_idx:
                ref_parts = split_text_by_sentences(ref_text, max_chars=300)
                if not ref_parts and ref_text:
                    ref_parts = [ref_text]

                for ref_part_idx_local, ref_part_content in enumerate(ref_parts, 1):
                    ref_slide = prs.slides.add_slide(content_slide_layout)
                    add_background_image(ref_slide, prs, 'C:/Users/Lyna Bouzouita/projet/background_type.jpg')
                    add_page_number(ref_slide, str(len(prs.slides)))
                    add_slide_border(ref_slide, RGBColor(0, 64, 128))
                    
                    title_shape_ref = ref_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                    title_frame_ref = title_shape_ref.text_frame
                    title_frame_ref.text = f"Référence de la réponse {idx}"
                    if len(ref_parts) > 1:
                        title_frame_ref.text += f" (Partie {ref_part_idx_local}/{len(ref_parts)})"
                    apply_text_style(title_frame_ref, font_name="Georgia", font_size=28,
                                  bold=True, color=RGBColor(0, 64, 128), alignment=PP_ALIGN.CENTER)
                    
                    # Pour les références, chaque puce sur une ligne
                    format_text_as_bullets(
                        ref_slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(3.5)).text_frame,
                        ref_part_content,
                        font_size=18,
                        line_spacing_val=1.2,  # Interligne plus serré pour les refs
                        space_after_pt=5,
                        space_before_pt=3,
                        is_continuation=False,
                        force_bullet_first_line=True
                    )

    # Nettoyer titre_doc pour l'utiliser dans le nom de fichier de sortie
    # Le titre_doc original est utilisé pour l'affichage sur la diapositive de titre.
    # Pour le nom de fichier, nous avons besoin d'une version nettoyée.
    
    # Étape 1: Remplacer les espaces par des tirets bas
    safe_filename_component = titre_doc.replace(' ', '_')
    
    # Étape 2: Supprimer/remplacer les autres caractères non valides dans les noms de fichiers.
    # Cette expression régulière conserve les caractères alphanumériques (lettres, chiffres, _), 
    # les points (.), et les tirets (-). Tous les autres caractères sont remplacés par un tiret bas.
    safe_filename_component = re.sub(r'[^\w\._-]', '_', safe_filename_component)
    
    # Étape 3: Remplacer les tirets bas multiples consécutifs par un seul (optionnel, pour l'esthétique)
    safe_filename_component = re.sub(r'__+', '_', safe_filename_component)
    
    # Étape 4: Supprimer les tirets bas en début/fin de chaîne (optionnel, pour l'esthétique)
    safe_filename_component = safe_filename_component.strip('_')
    
    # Étape 5: S'assurer que le composant du nom de fichier n'est pas vide après le nettoyage
    if not safe_filename_component:
        safe_filename_component = "document" # Valeur par défaut si le nettoyage résulte en une chaîne vide

    output_path = f"Analyse_Document_{safe_filename_component}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(output_path)
    return output_path

# --- STREAMLIT APP ---
def main():
    st.set_page_config(layout="wide")
    st.title("📄➡️📊 Due Diligence Assistant Pro")
    st.write("Téléversez vos documents (PDF/DOCX), posez vos questions, et obtenez une présentation PowerPoint synthétique.")

    col1, col2 = st.columns(2)

    with col1:
        st.subheader("1. Charger les Documents")
        uploaded_files = st.file_uploader("Choisissez un ou plusieurs fichiers", 
                                        type=["pdf", "docx"], 
                                        label_visibility="collapsed",
                                        accept_multiple_files=True)
        
        if uploaded_files:
            try:
                all_text = ""
                file_titles = []
                for uploaded_file in uploaded_files:
                    file_text = lire_document(uploaded_file)
                    all_text += f"\n\n--- DOCUMENT: {uploaded_file.name} ---\n\n{file_text}"
                    file_titles.append(uploaded_file.name)
                
                st.success(f"✅ {len(uploaded_files)} documents chargés !")
                with st.expander("Aperçu du texte extrait"):
                    st.text_area("Aperçu:", all_text[:1500], height=200, label_visibility="collapsed")
                st.session_state.texte = all_text
                st.session_state.titre_doc = " + ".join([os.path.splitext(name)[0] for name in file_titles])
                # Ajout du sélecteur de langue
                langue = st.selectbox(
                    "🌍 Choisissez la langue de génération :",
                    ("Français", "English")
                )
                st.session_state.langue = langue
                
                # SUPPRIMER cette ligne qui affiche le même aperçu
                # st.text_area("🔎 Aperçu rapide du document :", texte[:1000], height=250)

            except ValueError as e:
                st.error(e)
                st.session_state.texte = None
            except Exception as e:
                st.error(f"Erreur inattendue lors de la lecture du fichier : {e}")
                st.session_state.texte = None
        else:
            st.session_state.texte = None


    with col2:
        st.subheader("2. Poser des Questions")
        if "questions" not in st.session_state:
            st.session_state.questions = []

        with st.form("question_form"):
            new_question = st.text_input("Entrez une question sur le document :", key="new_q")
            submitted = st.form_submit_button("➕ Ajouter la question")
            if submitted and new_question.strip():
                st.session_state.questions.append(new_question.strip())
                st.rerun() # Rerun to update the list display immediately

        if st.session_state.questions:
            st.write("Questions à poser :")
            # Display questions with a delete button
            questions_to_keep = []
            for i, q in enumerate(st.session_state.questions):
                 q_col, btn_col = st.columns([0.9, 0.1])
                 q_col.write(f"{i+1}. {q}")
                 if btn_col.button("🗑️", key=f"del_{i}"):
                     # Don't add this question to the keep list
                     pass
                 else:
                     questions_to_keep.append(q)
            
            # If the list changed, update session state and rerun
            if len(questions_to_keep) != len(st.session_state.questions):
                 st.session_state.questions = questions_to_keep
                 st.rerun()


    st.divider()
    st.subheader("3. Générer la Présentation")

    # Disable button if no text or no questions
    button_disabled = not st.session_state.get("texte") or not st.session_state.get("questions")
    
    if st.button("🚀 Générer le PowerPoint", disabled=button_disabled, type="primary"):
        if st.session_state.texte and st.session_state.questions:
            with st.spinner('Création de la présentation en cours...'):
                try:
                    resume, qa_pairs = analyser_document(
                        st.session_state.texte, 
                        st.session_state.questions,
                        st.session_state.langue
                    )
                    st.write("Analyse terminée. Création du fichier PowerPoint...")
                    st.subheader("🧠 Aperçu des réponses générées :")
                    for i, (question, reponse) in enumerate(qa_pairs, 1):
                        extrait = reponse[:150] + ("..." if len(reponse) > 150 else "")
                        st.markdown(f"**{i}.** ❓ *{question}*")
                        st.markdown(f"➡️ {extrait}")
                    
                    pptx_path = creer_presentation(resume, qa_pairs, st.session_state.titre_doc)
                    if pptx_path and os.path.exists(pptx_path):
                        st.success("✅ Présentation générée avec succès !")
                        with open(pptx_path, "rb") as f:
                            st.download_button(
                                label="⬇️ Télécharger le PowerPoint",
                                data=f,
                                file_name=os.path.basename(pptx_path), # Use only filename
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        # Clean up generated file? Optional.
                        # os.remove(pptx_path)
                    else:
                         st.error("La génération du fichier PowerPoint a échoué. Vérifiez les logs ou le template.")

                except Exception as e:
                    st.error(f"Une erreur est survenue lors de la génération : {e}")
                    print(f"Erreur Génération: {e}") # Log error to console too
        elif not st.session_state.get("texte"):
             st.warning("Veuillez d'abord charger un document.")
        else:
            st.warning("Veuillez ajouter au moins une question.")

    # Add some instructions or info
    st.sidebar.info(
        """
        **Comment utiliser :**
        1.  Chargez un fichier PDF ou DOCX.
        2.  Ajoutez les questions spécifiques que vous souhaitez poser sur le document.
        3.  Cliquez sur 'Générer le PowerPoint'.
        4.  Téléchargez la présentation générée.

        **Template PowerPoint :**
        L'application utilise un template (`exemple.pptx`). Assurez-vous que ce fichier existe et contient au moins 4 slides avec les layouts attendus :
        - Slide 0: Pour le Titre
        - Slide 1: Pour le Résumé
        - Slide 2: Modèle pour les Questions
        - Slide 3: Modèle pour les Réponses
        """
    )


if __name__ == "__main__":
    main()





